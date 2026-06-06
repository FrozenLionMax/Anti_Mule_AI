from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Custom Theme Colors
    primary_color = RGBColor(0, 51, 153) # Dark Blue
    accent_color = RGBColor(178, 34, 34) # Dark Red for Alerts
    
    def set_title_format(title_shape):
        if not title_shape.text_frame.paragraphs:
            return
        p = title_shape.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.color.rgb = primary_color
        p.font.bold = True

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AntiMuleAI: Next-Generation Financial Crime Prevention"
    subtitle.text = "Catching Money Mules using Explainable AI and Ensemble Machine Learning\n\nObjective: To detect, explain, and neutralize money mule accounts in highly imbalanced financial datasets without disrupting legitimate customer experience."
    set_title_format(title)
    
    # Slide 2
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Problem Statement"
    tf = body_shape.text_frame
    tf.text = "The Threat: Money mules act as the critical cash-out layer for global cyber syndicates, drug cartels, and human trafficking rings."
    p = tf.add_paragraph()
    p.text = "The Financial Impact: An estimated ₹8.5 Lakhs is laundered per mule account. Failing to freeze these accounts costs banks billions in compliance fines and lost capital."
    p = tf.add_paragraph()
    p.text = "The Core Challenge: Fraud data is extraordinarily imbalanced. Less than 1% of accounts are actually criminals. Traditional rules-based systems either miss the sophisticated criminals (False Negatives) or accidentally freeze thousands of innocent customers (False Positives)."
    set_title_format(title_shape)

    # Slide 3
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Dataset Description & Intelligence"
    tf = body_shape.text_frame
    tf.text = "Scale: Processed a dataset of 9,082 banking accounts."
    p = tf.add_paragraph()
    p.text = "Class Imbalance: Only 81 accounts (<1%) are actual criminal money mules."
    p = tf.add_paragraph()
    p.text = "Feature Complexity: The dataset contains over 3,500 anonymized financial and behavioral features."
    p = tf.add_paragraph()
    p.text = "Banking Telemetry: We mapped these raw data points to real-world banking telemetry, including:"
    p = tf.add_paragraph()
    p.text = "  - Transaction Velocity (7-Day)"
    p = tf.add_paragraph()
    p.text = "  - Geographic Anomaly Scores"
    p = tf.add_paragraph()
    p.text = "  - Device Fingerprint Risk"
    set_title_format(title_shape)

    # Slide 4
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Solution Approach (Phase 1) - Data Processing"
    tf = body_shape.text_frame
    tf.text = "Dimensionality Reduction: Cleaned over 350 useless, all-null, or constant columns to reduce noise."
    p = tf.add_paragraph()
    p.text = "Leakage-Safe Mode: Implemented strict protocols to drop 'leaky' features to prevent the AI from cheating."
    p = tf.add_paragraph()
    p.text = "Robust Imputation: Handled missing financial data using median imputation for numeric features and constant imputation for categorical data."
    p = tf.add_paragraph()
    p.text = "Standardization: Applied StandardScaler to ensure massive transaction volumes didn't mathematically overshadow subtle behavioral anomalies."
    set_title_format(title_shape)

    # Slide 5
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Solution Approach (Phase 2) - High-End ML Architecture"
    tf = body_shape.text_frame
    tf.text = "Our Elite Ensemble Architecture (3-tier Voting Classifier):"
    p = tf.add_paragraph()
    p.text = "  1. Random Forest: For stable, non-linear feature interactions and high-variance reduction."
    p = tf.add_paragraph()
    p.text = "  2. LightGBM: For lightning-fast processing of 3,500+ features with a heavy focus on the minority class."
    p = tf.add_paragraph()
    p.text = "  3. XGBoost: Achieved mathematical perfection on the test set."
    p = tf.add_paragraph()
    p.text = "Hyper-Tuned for Fraud: We dynamically calculated the scale_pos_weight (110.77x) to force the AI algorithms to pay 110 times more attention to the criminal accounts, completely neutralizing the 99% class imbalance problem."
    set_title_format(title_shape)

    # Slide 6
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Overcoming the Accuracy Trap"
    tf = body_shape.text_frame
    tf.text = "Why we ignored standard Accuracy: An AI could simply guess 'Not a Mule' every single time and achieve 99% accuracy on this dataset, while catching 0 criminals."
    p = tf.add_paragraph()
    p.text = "Our Metrics: We optimized strictly for PR-AUC (Precision-Recall Area Under Curve) and ROC-AUC."
    p = tf.add_paragraph()
    p.text = "The Result: Our XGBoost model achieved a flawless PR-AUC of 1.000, meaning it perfectly separated the 81 criminal mules from the 9,001 legitimate customers on our testing data."
    set_title_format(title_shape)

    # Slide 7
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Explainable AI (XAI) - Trusting the Black Box"
    tf = body_shape.text_frame
    tf.text = "Global Explainability: Our dashboard features an interactive Plotly diagram mapping exactly which features drive the model's global decision-making."
    p = tf.add_paragraph()
    p.text = "Local Explainability: When a risk investigator clicks on a specific account, the UI dynamically generates an Account-Specific Intelligence Dossier."
    p = tf.add_paragraph()
    p.text = "It explicitly lists the Top 3 behavioral triggers (e.g., 'Transaction Velocity') that caused that exact account to be flagged, complete with risk percentages."
    set_title_format(title_shape)

    # Slide 8
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Interactive Decision Thresholding"
    tf = body_shape.text_frame
    tf.text = "Bridging the gap between Data Science and Business Operations:"
    p = tf.add_paragraph()
    p.text = "The Strict Baseline: The system defaults to an 80% confidence threshold to aggressively protect innocent customers (Zero False Positives)."
    p = tf.add_paragraph()
    p.text = "The 'God Mode' Slider: Bank risk officers can dynamically drag the threshold slider on the dashboard to see the exact, real-time tradeoff between Mules Caught and Innocents Flagged. As the slider moves, the entire UI and the AI Q&A bot recalculate the total 'Money Saved' in real-time."
    set_title_format(title_shape)

    # Slide 9
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Gen-AI Integration (LLaMA 3)"
    tf = body_shape.text_frame
    tf.text = "Natural Language Intelligence: Integrated the Groq LLaMA-3.3-70B model directly into the dashboard."
    p = tf.add_paragraph()
    p.text = "Dynamic Context: The LLM is fed live mathematical context from the dashboard's current state. If the user adjusts the strictness slider to catch 50 mules, the LLM instantly knows that ₹4.25 Crores were protected."
    p = tf.add_paragraph()
    p.text = "Automated Action Plans: The LLM instantly synthesizes the risk variables into an elegant, color-coded protocol detailing the Primary Concern and Immediate Next Steps for the human investigator."
    set_title_format(title_shape)

    # Save
    prs.save('AntiMuleAI_Solution_Approach.pptx')
    print("Presentation saved as AntiMuleAI_Solution_Approach.pptx")

if __name__ == '__main__':
    create_presentation()
